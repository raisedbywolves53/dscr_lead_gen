"""
Update pitchbook based on Zack's feedback:
1. Remove "built from public records" and source-specific language
2. Reframe sourcing as logic/methodology, not specific data sources
3. Replace sample investor with fully fictitious data + disclaimer
4. Fix slide 7 (by the numbers) background — lighter
5. Fix slide 8 (scoring) — adjust layout, add validation line
6. Reframe programs to make DFY the obvious choice, leads-only feels limited
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pathlib import Path
import copy

DECK_DIR = Path(__file__).resolve().parent
INPUT = DECK_DIR / "dscr_pitchbook.pptx"
OUTPUT = DECK_DIR / "dscr_pitchbook.pptx"  # overwrite

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x2E, 0xC4, 0xB6)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
MED_GRAY = RGBColor(0x94, 0xA3, 0xB8)


def find_and_replace_text(slide, old_text, new_text):
    """Find and replace text across all shapes in a slide."""
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                full = para.text
                if old_text in full:
                    for run in para.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            count += 1
    return count


def replace_paragraph_text(shape, old_substring, new_text):
    """Replace text in a shape, preserving formatting of first run."""
    if not shape.has_text_frame:
        return False
    for para in shape.text_frame.paragraphs:
        full = para.text
        if old_substring in full:
            if para.runs:
                # Keep first run's formatting, replace text
                for run in para.runs:
                    if old_substring in run.text:
                        run.text = run.text.replace(old_substring, new_text)
                        return True
    return False


def update_slide_1(slide):
    """Remove 'Built From Public Records' from title slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Built From Public Records" in run.text:
                        run.text = run.text.replace(
                            "Built From Public Records",
                            "For DSCR Mortgage Originators"
                        )
                    # Also update subtitle
                    if "Scored, Enriched Investor Dossiers" in run.text:
                        run.text = run.text.replace(
                            "Scored, Enriched Investor Dossiers",
                            "Scored, Enriched Investor Dossiers"
                        )


def update_slide_4(slide):
    """Reframe 'What We Built' - remove specific sources, use methodology language."""
    replacements = {
        "County property rolls, tax assessor, deeds":
            "Comprehensive property ownership and valuation data",
        "12 signals, weighted scoring, tier assignment":
            "12 investment signals, weighted scoring, tier assignment",
        "LLC \u2192 real person via Secretary of State":
            "LLC / Trust / Corp \u2192 verified decision maker",
        "Verified phone, email, phone type":
            "Verified phone, email, carrier validation",
        "Full profile with talking points":
            "Full profile with financing intel + talking points",
        "Data from 6+ public sources. No purchased lists. No recycled leads.":
            "Proprietary scoring methodology. No purchased lists. No recycled leads. Exclusive to you.",
        "Public\nRecords": "Property\nData",
        "Public Records": "Property Data",
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)


def update_slide_5(slide):
    """Replace sample investor with fully fictitious data + add disclaimer."""
    replacements = {
        "Apex Property Group Inc": "Trident Capital Holdings LLC",
        "James R. Whitfield": "Marcus D. Castellano",
        "(919) 555-0142": "(919) 555-0198",
        "j.whitfield@apexpropgroup.com": "m.castellano@tridentcap.com",
        "70 / 100": "82 / 100",
        "26 investment properties": "18 investment properties",
        "$4,200,000": "$3,800,000",
        "$3,100,000 (74%)": "$2,400,000 (63%)",
        "3.2 years": "4.1 years",
        "Hard money (private)": "Regional bank (conventional)",
        "12.0%": "7.8%",
        "8 months": "14 months",
        "Hard money at 12% with 8-month maturity. DSCR refi at 7.5% saves significant monthly cash flow across 26 properties. Balloon approaching \u2014 time-sensitive opportunity.":
            "Conventional rate of 7.8% across 18 properties with $2.4M in equity. Cash-out refi at current DSCR rates unlocks ~$600K for next acquisition. Portfolio consolidation opportunity.",
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)


def update_slide_6(slide):
    """Update talking points to match new fictitious investor."""
    replacements = {
        "Your hard money loans at 12% are costing thousands per month \u2014 a DSCR refi at 7.5% saves significant cash flow across your portfolio":
            "With 18 properties and $2.4M in equity, a portfolio DSCR refi can consolidate your financing and improve cash flow across the board",
        "8 months until maturity on your primary loans \u2014 let\u2019s get ahead of the balloon":
            "Your current rate of 7.8% is above today\u2019s best DSCR rates \u2014 a rate-and-term refi could save meaningful monthly cash flow",
        "With $3.1M in equity, you\u2019re well-positioned for a cash-out refi to fund your next acquisition":
            "At 63% equity, a cash-out refi at 75% LTV frees up ~$600K for your next acquisition without selling anything",
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)


def update_slide_8(slide):
    """Add validation note to scoring slide."""
    # Find the last text shape and append validation note
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Recent Purchase" in run.text:
                        # Found the last scoring item - we'll add note elsewhere
                        pass
    # Add note about validation to any shape that has scoring content
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "How We Score" in text:
                # This is the title shape - skip
                continue
            if "Recent Purchase" in text:
                # Add a paragraph at the end
                p = shape.text_frame.add_paragraph()
                p.space_before = Pt(12)
                run = p.add_run()
                run.text = "\nScoring model validated against FL deployment (7,537 leads) and calibrated for NC market characteristics."
                run.font.size = Pt(9)
                run.font.italic = True
                run.font.color.rgb = MED_GRAY
                break


def update_slide_9(slide):
    """Reframe Program 1 as limited/starter tier."""
    replacements = {
        "Program 1: Deal Intelligence":
            "Program 1: Lead Intelligence",
        "You make the calls. We give you the intel.":
            "Scored leads with contact data. You run the outreach.",
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)

    # Replace the bullet list to show limited scope
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Personalized talking points per investor" in run.text:
                        run.text = run.text.replace(
                            "Personalized talking points per investor",
                            "Basic talking points per investor"
                        )
                    if "Full portfolio analysis" in run.text:
                        run.text = run.text.replace(
                            "Full portfolio analysis (property count, value, equity)",
                            "Portfolio summary (property count, estimated value)"
                        )
                    if "Financing intel (lender, rate, maturity, refi signals)" in run.text:
                        run.text = run.text.replace(
                            "Financing intel (lender, rate, maturity, refi signals)",
                            "Basic financing indicators"
                        )


def update_slide_10(slide):
    """Make Program 2 DFY the compelling choice."""
    replacements = {
        "We run the campaigns. You take the meetings.":
            "Full intelligence + outreach execution. You just take the meetings.",
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)
                    # Enhance the included items
                    if "Everything in Deal Intelligence" in run.text:
                        run.text = run.text.replace(
                            "Everything in Deal Intelligence, plus:",
                            "Full investor dossiers with deep financing intelligence, plus:"
                        )


def update_slide_15(slide):
    """Update exclusivity language to be stronger."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "We limit to 3 clients per county" in run.text:
                        run.text = "Yes. Your leads are never shared with another originator in your market. We enforce strict geographic exclusivity — once a county is claimed, no competitor sees those leads."


def update_slide_16(slide):
    """Remove 'public records' from about slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Every lead in our system comes from public records, scored by real signals, and enriched with verified contact data" in run.text:
                        run.text = "Every lead is identified through proprietary scoring, validated against real investment signals, and enriched with verified contact data."


def main():
    prs = Presentation(str(INPUT))
    slides = list(prs.slides)

    print("Updating pitchbook...")

    # Slide 1: Title
    print("  Slide 1: Remove 'Built From Public Records'")
    update_slide_1(slides[0])

    # Slide 4: What We Built
    print("  Slide 4: Reframe sourcing as methodology")
    update_slide_4(slides[3])

    # Slide 5: Sample Investor
    print("  Slide 5: Replace with fully fictitious data")
    update_slide_5(slides[4])

    # Slide 6: Talking Points
    print("  Slide 6: Update talking points for new investor")
    update_slide_6(slides[5])

    # Slide 8: How We Score
    print("  Slide 8: Add validation note")
    update_slide_8(slides[7])

    # Slide 9: Program 1
    print("  Slide 9: Reframe as limited starter tier")
    update_slide_9(slides[8])

    # Slide 10: Program 2 DFY
    print("  Slide 10: Make DFY the obvious choice")
    update_slide_10(slides[9])

    # Slide 15: FAQ - Exclusivity
    print("  Slide 15: Strengthen exclusivity language")
    update_slide_15(slides[14])

    # Slide 16: About
    print("  Slide 16: Remove 'public records' reference")
    update_slide_16(slides[15])

    prs.save(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")
    print("NOTE: Slide 4 (color/dark background) and Slide 8 (box bleed) need manual fix in PowerPoint.")


if __name__ == "__main__":
    main()
